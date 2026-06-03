import io
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import requests
import os

def create_pptx(slides: List[Dict], title: str = "Research Paper Summary") -> io.BytesIO:
    """
    Generate a .pptx file from slide data.
    Returns a BytesIO stream containing the PowerPoint file.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    DARK_BG = RGBColor(26, 26, 46)
    DARK_CARD = RGBColor(22, 33, 62)
    ACCENT_BLUE = RGBColor(59, 130, 246)
    ACCENT_PURPLE = RGBColor(139, 92, 246)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(148, 163, 184)
    BORDER_COLOR = RGBColor(40, 40, 60)

    for idx, slide_data in enumerate(slides):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set slide background to dark
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

        # Add top accent bar
        top_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            prs.slide_width, Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ACCENT_BLUE
        top_bar.line.fill.background()

        # Slide number badge
        slide_num_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(2), Inches(0.5)
        )
        tf = slide_num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"SLIDE {idx + 1}"
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_BLUE
        p.font.bold = True

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.2),
            Inches(11.5), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", f"Slide {idx + 1}")
        p.font.size = Pt(36)
        p.font.color.rgb = WHITE
        p.font.bold = True

        # Accent line under title
        accent_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(2.4),
            Inches(2), Inches(0.06)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_PURPLE
        accent_line.line.fill.background()

        # Bullet points
        bullets = slide_data.get("bullets", [])
        bullet_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.8),
            Inches(11.5), Inches(4.0)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"✦  {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = LIGHT_GRAY
            p.space_after = Pt(12)
            p.level = 0

            # Alternate text color for visual hierarchy
            if i % 2 == 0:
                p.font.color.rgb = RGBColor(220, 220, 240)
            else:
                p.font.color.rgb = LIGHT_GRAY

        # Bottom accent bar
        bottom_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(7.35),
            prs.slide_width, Inches(0.15)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = ACCENT_PURPLE
        bottom_bar.line.fill.background()

        # Footer text
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),
            Inches(6), Inches(0.5)
        )
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = "AI Research Paper Summarizer"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 110, 130)

        # Page number
        page_num_box = slide.shapes.add_textbox(
            Inches(11), Inches(6.8),
            Inches(2), Inches(0.5)
        )
        tf = page_num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{idx + 1} / {len(slides)}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 110, 130)
        p.alignment = PP_ALIGN.RIGHT

    # Save to BytesIO
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer